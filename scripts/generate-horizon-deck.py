#!/usr/bin/env python3
"""
Horizon Scan PPTX Generator
============================
Generates a styled PPTX deck from structured JSON, using the CSA
AI x Security template as the visual master.

Architecture:
  RAG pipeline  →  JSON (pptx-schema.json)  →  this script  →  .pptx

The RAG system owns content only. This script owns all PPTX rendering.
No formatting logic should live in the JSON input.

Usage:
  python3 scripts/generate-horizon-deck.py [input.json] [output.pptx]

Defaults:
  input  : data/sample-horizon-scan.json
  output : outputs/final/generated-horizon-scan.pptx
"""

import json
import sys
import copy
from pathlib import Path
from xml.sax.saxutils import escape as xml_escape

import lxml.etree as etree
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ── Paths ─────────────────────────────────────────────────────────────────────

ROOT = Path(__file__).parent.parent
TEMPLATE_PATH = ROOT / "templates" / "AI x Security (for AISP projection) (1).pptx"
DEFAULT_INPUT  = ROOT / "data" / "sample-horizon-scan.json"
DEFAULT_OUTPUT = ROOT / "outputs" / "final" / "generated-horizon-scan.pptx"

# ── Layout index constants ────────────────────────────────────────────────────

LAYOUT_COVER          = 0   # Cover Page:        idx=0 CENTER_TITLE, idx=1 SUBTITLE
LAYOUT_TITLE_CONTENT  = 1   # Title and Content: idx=0 TITLE, idx=1 OBJECT
LAYOUT_TWO_CONTENT    = 2   # Two Content:       idx=0 TITLE, idx=1 OBJECT, idx=2 OBJECT
LAYOUT_COMPARISON     = 3   # Comparison:        idx=0 TITLE, idx=1/3 BODY, idx=2/4 OBJECT
LAYOUT_DIVIDER        = 4   # Divider:           idx=0 TITLE (centred, large)
LAYOUT_BLANK          = 5   # Blank
LAYOUT_SECTION_HEADER = 6   # Section Header:    idx=0 TITLE + background image

# slide_type string → layout index
LAYOUT_MAP = {
    "cover":            LAYOUT_COVER,
    "section_divider":  LAYOUT_DIVIDER,
    "section_header":   LAYOUT_SECTION_HEADER,
    "bullets":          LAYOUT_TITLE_CONTENT,
    "two_column":       LAYOUT_TWO_CONTENT,
    "comparison":       LAYOUT_COMPARISON,
    "blank":            LAYOUT_BLANK,
    # Semantic aliases → Title and Content
    "executive_summary": LAYOUT_TITLE_CONTENT,
    "timeline":          LAYOUT_TITLE_CONTENT,
    "threat_category":   LAYOUT_TITLE_CONTENT,
    "ai_as_target":      LAYOUT_TITLE_CONTENT,
    "ai_as_threat":      LAYOUT_TITLE_CONTENT,
    "ai_as_tool":        LAYOUT_TITLE_CONTENT,
    "case_study":        LAYOUT_TITLE_CONTENT,
    "recommendations":   LAYOUT_TITLE_CONTENT,
    "sources":           LAYOUT_TITLE_CONTENT,
    "closing":           LAYOUT_TITLE_CONTENT,
}

# ── XML namespace shortcut ────────────────────────────────────────────────────

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"

def _a(tag): return f"{{{A}}}{tag}"
def _p(tag): return f"{{{P}}}{tag}"

# ── Slide clearing ────────────────────────────────────────────────────────────

def clear_slides(prs):
    """Remove all existing slides from the presentation, keeping master/layouts."""
    sldIdLst = prs.slides._sldIdLst
    for sld_id in list(sldIdLst):
        rId = sld_id.get(qn("r:id"))
        sldIdLst.remove(sld_id)
        # Remove from internal relationships dict
        prs.part._rels._rels.pop(rId, None)

# ── Text helpers ──────────────────────────────────────────────────────────────

def _normalise_items(content):
    """
    Normalise bullet content to a list of dicts:
      {"text": str, "level": int, "bold": bool}
    Accepts str, list[str], or list[dict].
    """
    if isinstance(content, str):
        return [{"text": content, "level": 0, "bold": False}]
    items = []
    for item in content:
        if isinstance(item, str):
            items.append({"text": item, "level": 0, "bold": False})
        elif isinstance(item, dict):
            items.append({
                "text":  item.get("text", ""),
                "level": item.get("level", 0),
                "bold":  item.get("bold", False),
            })
    return items


def _make_para_xml(text, level=0, bold=False):
    """Build an <a:p> XML element with a single text run."""
    safe = xml_escape(str(text))
    bold_attr = ' b="1"' if bold else ""
    xml = (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr lvl="{level}"/>'
        f'<a:r><a:rPr{bold_attr}/><a:t>{safe}</a:t></a:r>'
        f'</a:p>'
    )
    return etree.fromstring(xml)


def set_title(slide, text):
    """Set the title placeholder (idx=0) text, preserving layout formatting."""
    try:
        ph = slide.placeholders[0]
    except KeyError:
        return
    tf = ph.text_frame
    txBody = tf._txBody
    # Remove all existing <a:p> elements
    for p_el in list(txBody.findall(_a("p"))):
        txBody.remove(p_el)
    # Add single paragraph
    txBody.append(_make_para_xml(text, level=0, bold=False))


def set_content_bullets(slide, ph_idx, items):
    """
    Fill a content placeholder with bullet paragraphs.
    items: list of {"text", "level", "bold"} dicts.
    """
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return
    tf = ph.text_frame
    tf.word_wrap = True
    txBody = tf._txBody
    # Remove existing paragraphs
    for p_el in list(txBody.findall(_a("p"))):
        txBody.remove(p_el)
    for item in items:
        txBody.append(_make_para_xml(item["text"], item["level"], item["bold"]))


def append_citations(slide, citations, ph_idx=1):
    """
    Append citation lines as small italic paragraphs to the content placeholder.
    Citations are rendered at level 4 with italics for visual separation.
    """
    if not citations:
        return
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return
    tf = ph.text_frame
    txBody = tf._txBody
    # Blank separator line
    sep_xml = f'<a:p xmlns:a="{A}"><a:pPr lvl="4"/><a:r><a:t></a:t></a:r></a:p>'
    txBody.append(etree.fromstring(sep_xml))
    for cite in citations:
        safe = xml_escape(str(cite))
        cit_xml = (
            f'<a:p xmlns:a="{A}">'
            f'<a:pPr lvl="4"/>'
            f'<a:r><a:rPr i="1"/><a:t>Source: {safe}</a:t></a:r>'
            f'</a:p>'
        )
        txBody.append(etree.fromstring(cit_xml))


def set_notes(slide, text):
    """Add speaker notes text to a slide."""
    if not text:
        return
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = str(text)
    except Exception:
        pass

# ── Slide builders ────────────────────────────────────────────────────────────

def build_cover(prs, data):
    layout = prs.slide_layouts[LAYOUT_COVER]
    slide = prs.slides.add_slide(layout)

    set_title(slide, data.get("title", ""))

    subtitle_parts = []
    if data.get("subtitle"):   subtitle_parts.append(data["subtitle"])
    if data.get("date"):       subtitle_parts.append(data["date"])
    if data.get("presenter"):  subtitle_parts.append(data["presenter"])
    if data.get("event"):      subtitle_parts.append(data["event"])
    subtitle = "\n".join(subtitle_parts)

    if subtitle:
        try:
            ph1 = slide.placeholders[1]
            tf = ph1.text_frame
            tf.word_wrap = True
            txBody = tf._txBody
            for p_el in list(txBody.findall(_a("p"))):
                txBody.remove(p_el)
            for i, line in enumerate(subtitle_parts):
                txBody.append(_make_para_xml(line, level=0, bold=(i == 0)))
        except KeyError:
            pass

    set_notes(slide, data.get("notes", ""))
    return slide


def build_divider(prs, data):
    layout = prs.slide_layouts[LAYOUT_DIVIDER]
    slide = prs.slides.add_slide(layout)
    set_title(slide, data.get("title", ""))
    if data.get("subtitle"):
        # Divider layout only has title; add subtitle as a text box if needed
        # For simplicity, append it to the title with a line break
        pass
    set_notes(slide, data.get("notes", ""))
    return slide


def build_section_header(prs, data):
    layout = prs.slide_layouts[LAYOUT_SECTION_HEADER]
    slide = prs.slides.add_slide(layout)
    set_title(slide, data.get("title", ""))
    set_notes(slide, data.get("notes", ""))
    return slide


def build_bullets(prs, data):
    layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]
    slide = prs.slides.add_slide(layout)
    set_title(slide, data.get("title", ""))

    raw_bullets = data.get("bullets", [])
    items = _normalise_items(raw_bullets)

    # Append key_stats as sub-bullets (level 1) if present
    stats = data.get("key_stats", [])
    if stats:
        items.append({"text": "", "level": 0, "bold": False})  # spacer
        for stat in stats:
            value = stat.get("value", "")
            label = stat.get("label", "")
            source = stat.get("source", "")
            line = f"{value}  —  {label}"
            if source:
                line += f"  ({source})"
            items.append({"text": line, "level": 1, "bold": True})

    # Append singapore_relevance if present
    sg = data.get("singapore_relevance", "")
    if sg:
        items.append({"text": "", "level": 0, "bold": False})
        items.append({"text": f"SG Relevance: {sg}", "level": 1, "bold": False})

    set_content_bullets(slide, 1, items)
    append_citations(slide, data.get("citations", []), ph_idx=1)
    set_notes(slide, data.get("notes", ""))
    return slide


def build_two_column(prs, data):
    layout = prs.slide_layouts[LAYOUT_TWO_CONTENT]
    slide = prs.slides.add_slide(layout)
    set_title(slide, data.get("title", ""))

    def col_items(col_data):
        items = []
        if col_data.get("header"):
            items.append({"text": col_data["header"], "level": 0, "bold": True})
        for b in col_data.get("bullets", []):
            items.extend(_normalise_items([b]))
        return items

    left_items  = col_items(data.get("left_column",  {}))
    right_items = col_items(data.get("right_column", {}))

    set_content_bullets(slide, 1, left_items)
    set_content_bullets(slide, 2, right_items)
    append_citations(slide, data.get("citations", []), ph_idx=1)
    set_notes(slide, data.get("notes", ""))
    return slide


def build_comparison(prs, data):
    layout = prs.slide_layouts[LAYOUT_COMPARISON]
    slide = prs.slides.add_slide(layout)
    set_title(slide, data.get("title", ""))

    # Comparison layout: idx=1 left header, idx=2 left content,
    #                    idx=3 right header, idx=4 right content
    left_header  = data.get("left_header",  "")
    right_header = data.get("right_header", "")
    left_bullets  = _normalise_items(data.get("left_bullets",  []))
    right_bullets = _normalise_items(data.get("right_bullets", []))

    if left_header:
        try:
            ph1 = slide.placeholders[1]
            tf = ph1.text_frame
            txBody = tf._txBody
            for p_el in list(txBody.findall(_a("p"))):
                txBody.remove(p_el)
            txBody.append(_make_para_xml(left_header, level=0, bold=True))
        except KeyError:
            pass

    if right_header:
        try:
            ph3 = slide.placeholders[3]
            tf = ph3.text_frame
            txBody = tf._txBody
            for p_el in list(txBody.findall(_a("p"))):
                txBody.remove(p_el)
            txBody.append(_make_para_xml(right_header, level=0, bold=True))
        except KeyError:
            pass

    set_content_bullets(slide, 2, left_bullets)
    set_content_bullets(slide, 4, right_bullets)
    append_citations(slide, data.get("citations", []), ph_idx=2)
    set_notes(slide, data.get("notes", ""))
    return slide


def build_blank(prs, data):
    layout = prs.slide_layouts[LAYOUT_BLANK]
    slide = prs.slides.add_slide(layout)
    set_notes(slide, data.get("notes", ""))
    return slide


# ── Dispatcher ────────────────────────────────────────────────────────────────

BUILDERS = {
    "cover":            build_cover,
    "section_divider":  build_divider,
    "section_header":   build_section_header,
    "two_column":       build_two_column,
    "comparison":       build_comparison,
    "blank":            build_blank,
}

def build_slide(prs, data):
    """Dispatch to the correct builder based on slide_type."""
    slide_type = data.get("slide_type", "bullets")
    builder = BUILDERS.get(slide_type, build_bullets)
    return builder(prs, data)

# ── Main ──────────────────────────────────────────────────────────────────────

def generate(input_path, output_path):
    input_path  = Path(input_path)
    output_path = Path(output_path)

    print(f"\n{'='*60}")
    print(f"  Horizon Scan PPTX Generator")
    print(f"{'='*60}")
    print(f"  Template : {TEMPLATE_PATH.relative_to(ROOT)}")
    print(f"  Input    : {input_path.relative_to(ROOT)}")
    print(f"  Output   : {output_path.relative_to(ROOT)}\n")

    # Load JSON
    deck_data = json.loads(input_path.read_text())
    slides_data = deck_data.get("slides", [])
    print(f"  Slides to generate: {len(slides_data)}")
    print(f"  Period: {deck_data.get('period', '—')}")
    print(f"  Sources analysed: {deck_data.get('source_count', '—')}\n")

    # Load template and clear existing slides
    print("  Loading template...")
    prs = Presentation(str(TEMPLATE_PATH))
    print(f"  Template has {len(prs.slides)} slides, {len(prs.slide_layouts)} layouts")
    print("  Clearing template slides...")
    clear_slides(prs)
    print(f"  Cleared. Slides remaining: {len(prs.slides)}")
    print()

    # Build slides
    type_counts = {}
    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("slide_type", "bullets")
        title = slide_data.get("title", "")[:50]
        slide = build_slide(prs, slide_data)
        type_counts[slide_type] = type_counts.get(slide_type, 0) + 1
        print(f"  [{i+1:2d}] {slide_type:<20} │ {title}")

    print(f"\n  ── Slide type summary ───────────────────")
    for stype, count in sorted(type_counts.items()):
        print(f"  {stype:<24} {count}")

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    size_kb = output_path.stat().st_size // 1024
    print(f"\n  ✓ Saved {len(slides_data)} slides → {output_path.relative_to(ROOT)} ({size_kb} KB)\n")
    return output_path


def main():
    input_path  = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_INPUT
    output_path = sys.argv[2] if len(sys.argv) > 2 else DEFAULT_OUTPUT
    generate(input_path, output_path)


if __name__ == "__main__":
    main()
