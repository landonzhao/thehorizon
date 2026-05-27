#!/usr/bin/env python3
"""
PPTX Template Inspector
Enumerates all slide layouts, placeholders, text boxes, theme colors,
and shape geometry from a .pptx template file.
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "AI x Security (for AISP projection) (1).pptx"


def rgb_to_hex(rgb):
    if rgb is None:
        return None
    try:
        return "#{:02X}{:02X}{:02X}".format(rgb.r, rgb.g, rgb.b)
    except Exception:
        return str(rgb)


def emu_to_inches(emu):
    return round(emu / 914400, 3)


def inspect_placeholder(ph):
    info = {
        "idx": ph.placeholder_format.idx,
        "type": str(ph.placeholder_format.type),
        "name": ph.name,
        "left": emu_to_inches(ph.left) if ph.left is not None else None,
        "top": emu_to_inches(ph.top) if ph.top is not None else None,
        "width": emu_to_inches(ph.width) if ph.width is not None else None,
        "height": emu_to_inches(ph.height) if ph.height is not None else None,
    }
    # Sample text if present
    try:
        if ph.has_text_frame and ph.text_frame.text:
            info["sample_text"] = ph.text_frame.text[:80]
    except Exception:
        pass
    # Font info from first paragraph/run
    try:
        para = ph.text_frame.paragraphs[0]
        if para.runs:
            run = para.runs[0]
            info["font_name"] = run.font.name
            info["font_size"] = str(run.font.size.pt) if run.font.size else None
            info["font_bold"] = run.font.bold
            info["font_color"] = rgb_to_hex(run.font.color.rgb) if run.font.color and run.font.color.type else None
    except Exception:
        pass
    return info


def inspect_shape(shape):
    info = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": emu_to_inches(shape.left) if shape.left is not None else None,
        "top": emu_to_inches(shape.top) if shape.top is not None else None,
        "width": emu_to_inches(shape.width) if shape.width is not None else None,
        "height": emu_to_inches(shape.height) if shape.height is not None else None,
    }
    if shape.has_text_frame:
        info["text"] = shape.text_frame.text[:120]
        paras = []
        for para in shape.text_frame.paragraphs:
            p = {"text": para.text[:80]}
            try:
                p["alignment"] = str(para.alignment)
            except Exception:
                pass
            if para.runs:
                r = para.runs[0]
                p["font_name"] = r.font.name
                p["font_size"] = str(r.font.size.pt) if r.font.size else None
                p["bold"] = r.font.bold
                p["color"] = rgb_to_hex(r.font.color.rgb) if r.font.color and r.font.color.type else None
            paras.append(p)
        info["paragraphs"] = paras
    # Fill color
    try:
        fill = shape.fill
        if fill.type is not None:
            if str(fill.type) == "SOLID (1)":
                info["fill_color"] = rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return info


def inspect_layout(layout, idx):
    info = {
        "index": idx,
        "name": layout.name,
        "placeholders": [],
        "shapes": [],
    }
    for ph in layout.placeholders:
        info["placeholders"].append(inspect_placeholder(ph))
    for shape in layout.shapes:
        try:
            _ = shape.placeholder_format
            is_ph = True
        except (ValueError, AttributeError):
            is_ph = False
        if not is_ph:
            info["shapes"].append(inspect_shape(shape))
    return info


def inspect_theme_colors(prs):
    colors = {}
    try:
        theme = prs.slide_master.theme_color_map
        for name, color in theme.items():
            colors[name] = rgb_to_hex(color)
    except Exception:
        pass
    # Try reading directly from XML
    try:
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        theme_el = prs.slide_master.element.find(".//" + qn("a:theme"))
        if theme_el is not None:
            fmtScheme = theme_el.find(".//" + qn("a:fmtScheme"))
            clrScheme = theme_el.find(".//" + qn("a:clrScheme"))
            if clrScheme is not None:
                for child in clrScheme:
                    tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                    for rgb_el in child:
                        val = rgb_el.get("val") or rgb_el.get("lastClr") or rgb_el.get("lumMod")
                        if val:
                            colors[tag] = "#" + val.upper() if not val.startswith("#") else val.upper()
    except Exception:
        pass
    return colors


def inspect_slide_master_backgrounds(prs):
    results = []
    master = prs.slide_master
    for shape in master.shapes:
        s = inspect_shape(shape)
        if s.get("fill_color") or s.get("text"):
            results.append(s)
    return results


def main():
    path = TEMPLATE_PATH
    if len(sys.argv) > 1:
        path = Path(sys.argv[1])

    print(f"\n{'='*70}")
    print(f"  PPTX Template Inspector")
    print(f"  File: {path}")
    print(f"{'='*70}\n")

    prs = Presentation(str(path))

    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)
    print(f"Slide dimensions: {slide_w}\" × {slide_h}\"  ({prs.slide_width} × {prs.slide_height} EMU)\n")

    # Theme colors
    print("── Theme Colors ──────────────────────────────────────────────────────")
    theme_colors = inspect_theme_colors(prs)
    for k, v in theme_colors.items():
        print(f"  {k:<20} {v}")

    print()

    # Slide layouts
    print("── Slide Layouts ─────────────────────────────────────────────────────")
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        info = inspect_layout(layout, i)
        layouts.append(info)
        ph_summary = ", ".join(
            f"idx={p['idx']} ({p['type'].split('(')[0].strip()}, '{p['name']}')"
            for p in info["placeholders"]
        )
        print(f"\n  [{i:2d}] {layout.name}")
        print(f"       Placeholders: {ph_summary or '(none)'}")
        if info["shapes"]:
            shape_names = [s["name"] for s in info["shapes"]]
            print(f"       Extra shapes: {', '.join(shape_names)}")

    print()

    # Existing slides in the template (if any)
    print("── Existing Slides in Template ───────────────────────────────────────")
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        texts = [s.text_frame.text[:50] for s in slide.shapes if s.has_text_frame and s.text_frame.text]
        print(f"  Slide {i+1}: layout='{layout_name}'")
        for t in texts[:3]:
            print(f"    › {t.strip()}")

    # Save full inventory as JSON
    output = {
        "slide_width_inches": slide_w,
        "slide_height_inches": slide_h,
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "theme_colors": theme_colors,
        "layouts": layouts,
    }
    out_path = Path(__file__).parent.parent / "outputs" / "debug" / "template_inventory.json"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(json.dumps(output, indent=2))
    print(f"\n✓ Full inventory saved to: {out_path.relative_to(Path(__file__).parent.parent)}\n")

    return output


if __name__ == "__main__":
    main()
